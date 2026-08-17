#!/usr/bin/env python
"""테스트베드 통과 기준 자동 점검 (11번 알고리즘설명서 ↔ 21번 리밸런싱 발생내역).

사용법:
    python verify_pass_criteria.py <11번_알고리즘설명서.pptx> <21번_리밸런싱발생내역.xlsx>

점검 기준 (SKILL.md '테스트베드 통과 기준' 참조):
  기준1 잔고정합: 각 그룹·리밸런싱일자(rday)별로 자산종류별
                  |실잔고비중 − 목표(MP)비중| 의 합계 < 20%
  기준2 위험도정합: 21번 실제 위험도 Σ(자산종류비중×위험도점수) 가
                    11번 위험도 범위 [최저,최고] 안
  기준2 자산비중정합: 21번 MP 자산종류별 실제 비중이 11번 자산종류별
                      편입비중표(b)의 허용범위 안

규칙·임계값(20%)은 코스콤 심사 고정 기준. 범위 기준은 11번 알고리즘설명서.
헤더/컬럼은 동적 탐지(하드코딩 없음). 그룹↔전략 매핑: 적극=매운맛, 중립=중간맛, 안정=순한맛.
"""
import sys
import openpyxl
from pptx import Presentation
from collections import defaultdict

GROUP = [("적극", "매운맛"), ("중립", "중간맛"), ("안정", "순한맛")]
STRAT = ["매운맛", "중간맛", "순한맛"]
BAL_TOL = 20.0   # 기준1 임계(%) — 자산종류별 비중차 합계 상한
W_TOL = 0.05     # 비중 범위 허용오차(%p)
D_TOL = 0.01     # 위험도 허용오차


def is_cash(isin):
    s = str(isin)
    return s in ("CASH", "현금") or "CASH" in s


def prng(s):
    """'0~90%'→(0,90), '2.5%'→(2.5,2.5)."""
    s = str(s).replace("%", "").strip()
    if "~" in s:
        a, b = s.split("~")
        return float(a), float(b)
    return float(s), float(s)


def _find_strat_cols(t, scan=4):
    """매운맛/중간맛/순한맛 컬럼 인덱스 탐지. 헤더가 여러 행에 걸쳐도(전략구분/
    참여유형 병기 등) 상위 scan행을 훑어 찾는다."""
    nr, nc = len(t.rows), len(t.columns)
    gcol = {}
    for gn in STRAT:
        for c in range(nc):
            if any(t.cell(r, c).text.strip() == gn for r in range(min(scan, nr))):
                gcol[gn] = c
                break
    return gcol


def pptx_ranges(pptx_path):
    """11번에서 자산종류별 비중 허용범위(편입비중표)와 위험도 범위를 추출.
    헤더가 1행이든 여러 행(자산종류 | 비중(참여유형/전략구분) | 공격형… | 매운맛…)
    이든 전략 컬럼을 동적 탐지한다."""
    prs = Presentation(pptx_path)
    brange = {}   # 자산종류 -> {전략: (lo,hi)}
    drange = {}   # 전략 -> (lo,hi)
    for sl in prs.slides:
        for sh in sl.shapes:
            if not sh.has_table:
                continue
            t = sh.table
            nr = len(t.rows)
            gcol = _find_strat_cols(t)
            if len(gcol) < 3:
                continue
            first = " ".join(t.cell(r, 0).text for r in range(min(4, nr)))
            # 편입비중표: 첫 컬럼이 '자산종류'(참여유형 아님) + 데이터행 비중이 % 값
            if "자산" in first and "종류" in first:
                for r in range(nr):
                    at = t.cell(r, 0).text.strip()
                    if not at or at in ("자산종류", "자산 종류", "합계"):
                        continue
                    if "%" not in t.cell(r, gcol["매운맛"]).text:
                        continue
                    for gn, c in gcol.items():
                        v = t.cell(r, c).text.strip()
                        if v:
                            brange.setdefault(at, {})[gn] = prng(v)
            # 위험도 범위표: '위험도 범위' 행 + 전략 컬럼의 숫자 값
            for r in range(nr):
                if "위험도 범위" in t.cell(r, 0).text:
                    for gn, c in gcol.items():
                        v = t.cell(r, c).text.strip()
                        if v and any(ch.isdigit() for ch in v):
                            drange[gn] = prng(v)
    return brange, drange


def xlsx_data(xlsx_path):
    """21번에서 ISIN→자산종류·위험도점수, MP 자산종류별 비중, 잔고 평가금액을 추출."""
    wb = openpyxl.load_workbook(xlsx_path, data_only=False)
    uni = wb["투자유니버스"]
    uh = [uni.cell(1, c).value for c in range(1, uni.max_column + 1)]
    col = lambda n: uh.index(n) + 1
    isin2at, at2sc = {}, {}
    for r in range(2, uni.max_row + 1):
        iv = uni.cell(r, col("종목코드(ISIN코드)")).value
        if iv:
            a = uni.cell(r, col("자산종류")).value
            isin2at[iv] = a
            at2sc[a] = uni.cell(r, col("위험도 점수")).value
    at2sc["현금"] = 1

    # MP 자산종류별 비중: 전략 -> rday -> 자산종류 -> 비중
    mp = defaultdict(lambda: defaultdict(lambda: defaultdict(float)))
    for gs, gn in GROUP:
        ws = wb[f"MP내역({gs})"]
        for r in range(2, ws.max_row + 1):
            d = ws.cell(r, 1).value
            isin = ws.cell(r, 2).value
            w = ws.cell(r, 7).value
            if not isin or isin == "합계" or w in (None, ""):
                continue
            a = "현금" if is_cash(isin) else isin2at.get(isin)
            mp[gn][str(d)[:10]][a] += float(w)

    # 잔고 평가금액: 전략 -> rday -> 자산종류 -> 평가금액 (실잔고비중은 합으로 정규화)
    bal = defaultdict(lambda: defaultdict(lambda: defaultdict(float)))
    for sn in wb.sheetnames:
        if not sn.startswith("잔고변경현황"):
            continue
        gn = next((g for gs, g in GROUP if gs in sn), None)
        ws = wb[sn]
        hdr = {ws.cell(4, c).value: c for c in range(1, ws.max_column + 1)}
        c_rday = hdr.get("리밸런싱일자")
        c_isin = hdr.get("ISIN코드")
        c_val = hdr.get("평가금액")
        if None in (c_rday, c_isin, c_val):
            continue
        for r in range(5, ws.max_row + 1):
            isin = ws.cell(r, c_isin).value
            if not isin or isin == "합계":
                continue
            rday = str(ws.cell(r, c_rday).value)[:10]
            val = ws.cell(r, c_val).value or 0
            a = "현금" if is_cash(isin) else isin2at.get(isin)
            bal[gn][rday][a] += float(val)
    return at2sc, mp, bal


def main(pptx_path, xlsx_path):
    brange, drange = pptx_ranges(pptx_path)
    at2sc, mp, bal = xlsx_data(xlsx_path)
    issues = []

    # 기준1 잔고정합: 자산종류별 |실잔고비중 − 목표비중| 합 < 20%
    for gn in mp:
        for rday, tgt in mp[gn].items():
            b = bal.get(gn, {}).get(rday)
            if not b:
                continue
            tot = sum(b.values())
            real = {a: v / tot for a, v in b.items()} if tot else {}
            ats = set(tgt) | set(real)
            diff = sum(abs(real.get(a, 0) - tgt.get(a, 0)) for a in ats) * 100
            if diff >= BAL_TOL:
                issues.append(f"[기준1 잔고정합] {gn} {rday}: 자산종류 비중차 합 {diff:.1f}% ≥ {BAL_TOL:.0f}%")

    # 기준2 위험도정합
    for gn in mp:
        for rday in mp[gn]:
            d = sum(w * at2sc.get(a, 1) for a, w in mp[gn][rday].items())
            rg = drange.get(gn)
            if not rg:
                continue
            lo, hi = rg
            if not (lo - D_TOL <= d <= hi + D_TOL):
                issues.append(f"[기준2 위험도] {gn} {rday}: 실측 {d:.3f} ∉ [{lo}, {hi}]")

    # 기준2 자산비중정합
    for gn in mp:
        for rday in mp[gn]:
            for a, w in mp[gn][rday].items():
                rg = brange.get(a, {}).get(gn)
                if not rg:
                    continue
                lo, hi = rg
                wp = w * 100
                if not (lo - W_TOL <= wp <= hi + W_TOL):
                    issues.append(f"[기준2 자산비중] {gn} {rday} {a}: 실측 {wp:.1f}% ∉ [{lo}, {hi}]")

    print(f"=== 테스트베드 통과 기준 점검 ===")
    print(f"11번: {pptx_path}")
    print(f"21번: {xlsx_path}")
    if not issues:
        print("✅ 전체 통과 — 기준1(잔고정합) · 기준2(위험도/자산비중 정합) 모두 충족")
        return 0
    print(f"⚠ 위반 {len(issues)}건:")
    for it in issues:
        print("  -", it)
    return 1


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("사용법: python verify_pass_criteria.py <11번.pptx> <21번.xlsx>")
        sys.exit(2)
    sys.exit(main(sys.argv[1], sys.argv[2]))
