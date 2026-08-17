"""pptx에서 모든 shape를 추출하여 shapes_manifest.json을 생성.

슬라이드별 텍스트박스, 테이블, 이미지, 그룹 등 전체 shape의
물리적 정보(위치, 크기, 내용, 서식)를 빠짐없이 추출한다.

Usage:
    python extract_all_shapes.py <input.pptx> [output.json]
"""
import sys, json
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE
from copy import deepcopy
from lxml import etree


# ── 테이블 스타일 추출 (첫 번째 테이블에서 한 번만) ──

def _extract_tcPr_xml(tc_element):
    """tcPr XML을 문자열로 직렬화."""
    tcPr = tc_element.find(qn('a:tcPr'))
    if tcPr is None:
        return None
    return etree.tostring(tcPr, encoding='unicode')


def extract_table_style(prs):
    """프레젠테이션에서 테이블 스타일 정보를 추출."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl_el = shape.table._tbl
            tblPr = tbl_el.find(qn('a:tblPr'))
            if tblPr is None:
                continue

            style = {
                "tblPr_xml": etree.tostring(tblPr, encoding='unicode'),
            }

            rows = tbl_el.findall(qn('a:tr'))
            if rows:
                tcs = rows[0].findall(qn('a:tc'))
                if len(tcs) >= 1:
                    style["header_tcPr_xml"] = _extract_tcPr_xml(tcs[0])
                if len(tcs) >= 2:
                    style["value_tcPr_xml"] = _extract_tcPr_xml(tcs[1])

            # value_tcPr 폴백: header와 동일하게
            if "value_tcPr_xml" not in style:
                style["value_tcPr_xml"] = style.get("header_tcPr_xml")

            # 빨간 셀 찾기 (FF0000)
            for s in prs.slides:
                for sh in s.shapes:
                    if not sh.has_table:
                        continue
                    for tr in sh.table._tbl.findall(qn('a:tr')):
                        for tc in tr.findall(qn('a:tc')):
                            tcPr = tc.find(qn('a:tcPr'))
                            if tcPr is None:
                                continue
                            fill = tcPr.find(qn('a:solidFill'))
                            if fill is None:
                                continue
                            srgb = fill.find(qn('a:srgbClr'))
                            if srgb is not None and srgb.get('val', '').upper() == 'FF0000':
                                style["red_tcPr_xml"] = _extract_tcPr_xml(tc)
                                return style

            # 빨간셀 없으면 value 스타일로 폴백
            if "red_tcPr_xml" not in style:
                style["red_tcPr_xml"] = style.get("value_tcPr_xml")

            return style
    return None


# ── 셀 스타일 판정 ──

def _cell_style(tc):
    """셀의 배경색으로 스타일(H/V/R) 판정."""
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        return "V"
    fill = tcPr.find(qn('a:solidFill'))
    if fill is None:
        return "V"
    srgb = fill.find(qn('a:srgbClr'))
    if srgb is not None and srgb.get('val', '').upper() == 'FF0000':
        return "R"
    scheme = fill.find(qn('a:schemeClr'))
    if scheme is not None:
        return "H"
    return "V"


def _cell_align(tc):
    """셀의 첫 paragraph 정렬."""
    txBody = tc.find(qn('a:txBody'))
    if txBody is None:
        return None
    p = txBody.find(qn('a:p'))
    if p is None:
        return None
    pPr = p.find(qn('a:pPr'))
    if pPr is None:
        return None
    algn = pPr.get('algn')
    return algn if algn and algn != 'ctr' else None


def _cell_merge(tc):
    """셀의 병합 정보."""
    merge = {}
    rs = tc.get('rowSpan')
    cs = tc.get('gridSpan')
    vm = tc.get('vMerge')
    hm = tc.get('hMerge')
    if rs:
        merge["rs"] = int(rs)
    if cs:
        merge["cs"] = int(cs)
    if vm:
        merge["vMerge"] = True
    if hm:
        merge["hMerge"] = True
    return merge or None


# ── 테이블 추출 ──

def extract_table(shape):
    """테이블 shape에서 구조 정보를 추출."""
    tbl = shape.table
    tbl_el = tbl._tbl

    # 열 너비
    grid = tbl_el.find(qn('a:tblGrid'))
    col_widths_cm = []
    for gc in grid.findall(qn('a:gridCol')):
        col_widths_cm.append(round(int(gc.get('w')) / 360000, 2))

    rows = []
    for ri, tr in enumerate(tbl_el.findall(qn('a:tr'))):
        h_cm = round(int(tr.get('h')) / 360000, 2)
        cells = []
        for ci, tc in enumerate(tr.findall(qn('a:tc'))):
            cell = {"v": tbl.cell(ri, ci).text.strip()}
            cell["s"] = _cell_style(tc)
            merge = _cell_merge(tc)
            if merge:
                cell["merge"] = merge
            algn = _cell_align(tc)
            if algn:
                cell["algn"] = algn
            cells.append(cell)
        rows.append({"cells": cells, "h": h_cm})

    return {
        "rows_count": len(tbl.rows),
        "cols_count": len(tbl.columns),
        "col_widths_cm": col_widths_cm,
        "rows": rows,
    }


# ── 텍스트 추출 ──

def _color_to_str(color):
    """ColorFormat을 문자열로 변환."""
    try:
        if color.type is not None:
            rgb = color.rgb
            if rgb:
                return str(rgb)
    except Exception:
        pass
    return None


def extract_paragraphs(text_frame):
    """TextFrame에서 paragraph별 서식 정보를 추출."""
    paragraphs = []
    for para in text_frame.paragraphs:
        text = para.text
        # paragraph 레벨 서식
        p_info = {"text": text}

        # 정렬
        if para.alignment is not None:
            align_map = {0: "l", 1: "ctr", 2: "r", 3: "justify"}
            p_info["align"] = align_map.get(int(para.alignment), "l")

        # 들여쓰기 레벨
        if para.level and para.level > 0:
            p_info["indent"] = para.level

        # run 레벨 서식 (첫 run 기준, 동일 paragraph 내 서식이 다르면 runs 배열로)
        runs = para.runs
        if runs:
            first = runs[0]
            if first.font.size is not None:
                p_info["font_size"] = round(first.font.size.pt, 1)
            if first.font.bold:
                p_info["bold"] = True
            if first.font.italic:
                p_info["italic"] = True
            color = _color_to_str(first.font.color)
            if color:
                p_info["color"] = color
            if first.font.name:
                p_info["font_name"] = first.font.name

            # 같은 paragraph에서 서식이 다른 run이 있는 경우
            if len(runs) > 1:
                has_mixed = False
                for r in runs[1:]:
                    if (r.font.bold != first.font.bold or
                        r.font.size != first.font.size or
                        r.font.italic != first.font.italic):
                        has_mixed = True
                        break
                if has_mixed:
                    p_info["runs"] = []
                    for r in runs:
                        run_info = {"text": r.text}
                        if r.font.size is not None:
                            run_info["font_size"] = round(r.font.size.pt, 1)
                        if r.font.bold:
                            run_info["bold"] = True
                        if r.font.italic:
                            run_info["italic"] = True
                        p_info["runs"].append(run_info)

        paragraphs.append(p_info)
    return paragraphs


def extract_textbox(shape):
    """텍스트박스 shape에서 내용과 서식을 추출."""
    tf = shape.text_frame
    return {
        "content": shape.text.strip(),
        "word_wrap": tf.word_wrap,
        "paragraphs": extract_paragraphs(tf),
    }


# ── Shape 추출 (재귀) ──

def extract_shape(shape, slide_idx, shape_counter):
    """단일 shape를 딕셔너리로 추출. 그룹이면 재귀."""
    shape_counter[0] += 1
    sid = f"s{slide_idx + 1}_shape{shape_counter[0]}"

    info = {
        "id": sid,
        "position": {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
        },
    }

    # shape 이름 (PowerPoint에서 부여한 이름)
    if shape.name:
        info["name"] = shape.name

    st = shape.shape_type
    if st is not None:
        info["shape_type"] = int(st)

    # 테이블
    if shape.has_table:
        info["type"] = "table"
        info["table"] = extract_table(shape)
        return info

    # 그룹
    if st == MSO_SHAPE_TYPE.GROUP:
        info["type"] = "group"
        children = []
        for child in shape.shapes:
            children.append(extract_shape(child, slide_idx, shape_counter))
        info["children"] = children
        return info

    # 이미지
    if st == MSO_SHAPE_TYPE.PICTURE:
        info["type"] = "image"
        try:
            img = shape.image
            info["image"] = {
                "content_type": img.content_type,
                "filename": img.filename if hasattr(img, 'filename') else None,
            }
        except Exception:
            info["image"] = {}
        return info

    # 텍스트가 있는 shape (텍스트박스, 도형 안 텍스트 등)
    if shape.has_text_frame:
        text = shape.text.strip()
        if text:
            info["type"] = "textbox"
            info["text"] = extract_textbox(shape)
            return info
        else:
            # 빈 텍스트 shape
            info["type"] = "empty_textbox"
            return info

    # 기타 (커넥터, 프리폼 등)
    info["type"] = "other"
    return info


# ── 레이아웃 추출 (pptx에서 확실히 추출 가능한 값만) ──

def extract_layout(prs, blank_layout_index=6, cell_margin_emu=36000):
    """pptx에서 프로그래매틱하게 확정할 수 있는 레이아웃 값만 추출.

    gap, heading_styles, margin 등은 레퍼런스 shape 구조와 출력 구조가 다르므로
    여기서 추출하지 않는다. Phase 2~4에서 AI가 스크린샷을 보고 결정.
    """
    from collections import Counter

    # font_name: pptx 테마 XML에서 minorFont ea 이름 해석
    font_name = None
    try:
        theme_el = prs.slide_masters[0].element.find(
            './/' + qn('a:minorFont'))
        if theme_el is not None:
            ea = theme_el.find(qn('a:ea'))
            if ea is not None:
                font_name = ea.get('typeface')
    except Exception:
        pass
    # 테마에서 못 찾으면 run에서 빈도 기반 추출
    if not font_name:
        names = Counter()
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and not run.font.name.startswith('+'):
                                names[run.font.name] += 1
        font_name = names.most_common(1)[0][0] if names else None

    # theme_font: 테이블 셀 run의 ea typeface (XML 레벨)
    theme_font = "+mn-ea"
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for tr in shape.table._tbl.findall(qn('a:tr')):
                for tc in tr.findall(qn('a:tc')):
                    txBody = tc.find(qn('a:txBody'))
                    if txBody is None:
                        continue
                    for p in txBody.findall(qn('a:p')):
                        for r in p.findall(qn('a:r')):
                            rPr = r.find(qn('a:rPr'))
                            if rPr is None:
                                continue
                            ea = rPr.find(qn('a:ea'))
                            if ea is not None and ea.get('typeface'):
                                theme_font = ea.get('typeface')
                                break

    # default_font_size: 테이블 셀에서 가장 빈번한 font size
    sz_counter = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for tr in shape.table._tbl.findall(qn('a:tr')):
                for tc in tr.findall(qn('a:tc')):
                    txBody = tc.find(qn('a:txBody'))
                    if txBody is None:
                        continue
                    for p in txBody.findall(qn('a:p')):
                        for r in p.findall(qn('a:r')):
                            rPr = r.find(qn('a:rPr'))
                            if rPr is not None and rPr.get('sz'):
                                sz_counter[rPr.get('sz')] += 1
    default_font_size = sz_counter.most_common(1)[0][0] if sz_counter else "1000"

    # cell_margin: 첫 번째 테이블 셀의 marL
    cell_margin = str(cell_margin_emu)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for tr in shape.table._tbl.findall(qn('a:tr')):
                for tc in tr.findall(qn('a:tc')):
                    tcPr = tc.find(qn('a:tcPr'))
                    if tcPr is not None and tcPr.get('marL'):
                        cell_margin = tcPr.get('marL')
                        break

    # blank_slide_layout_index: placeholder가 없는 레이아웃
    blank_idx = blank_layout_index
    for i, layout in enumerate(prs.slide_layouts):
        if len(layout.placeholders) == 0:
            blank_idx = i
            break

    return {
        "font_name": font_name,
        "theme_font": theme_font,
        "default_font_size_hundredths": default_font_size,
        "cell_margin_emu": cell_margin,
        "blank_slide_layout_index": blank_idx,
    }


# ── 메인 추출 ──

def extract_all(pptx_path, blank_layout_index=6, cell_margin_emu=36000):
    """pptx 파일에서 전체 shape manifest를 추출."""
    prs = Presentation(pptx_path)

    manifest = {
        "source": pptx_path,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slide_count": len(prs.slides),
        "layout": extract_layout(prs, blank_layout_index=blank_layout_index,
                                 cell_margin_emu=cell_margin_emu),
        "table_style": extract_table_style(prs),
        "slides": [],
    }

    total_shapes = 0
    for si, slide in enumerate(prs.slides):
        shape_counter = [0]
        shapes = []
        for shape in slide.shapes:
            shapes.append(extract_shape(shape, si, shape_counter))
        total_shapes += shape_counter[0]

        manifest["slides"].append({
            "slide_no": si + 1,
            "shape_count": shape_counter[0],
            "shapes": shapes,
        })

    manifest["total_shapes"] = total_shapes
    return manifest


# ── 요약 출력 ──

def print_summary(manifest):
    """추출 결과 요약을 stderr로 출력."""
    print(f"\n=== {manifest['source']} ===", file=sys.stderr)
    print(f"슬라이드: {manifest['slide_count']}장, "
          f"총 shape: {manifest['total_shapes']}개, "
          f"크기: {manifest['slide_width']}x{manifest['slide_height']} EMU",
          file=sys.stderr)

    for sl in manifest['slides']:
        types = {}
        for sh in sl['shapes']:
            t = sh['type']
            types[t] = types.get(t, 0) + 1

        parts = []
        for t in ['textbox', 'table', 'image', 'group', 'empty_textbox', 'other']:
            if t in types:
                parts.append(f"{t}={types[t]}")

        # 테이블 상세
        table_details = []
        for sh in sl['shapes']:
            if sh['type'] == 'table':
                tbl = sh['table']
                merges = sum(
                    1 for r in tbl['rows'] for c in r['cells']
                    if c.get('merge') and ('rs' in c['merge'] or 'cs' in c['merge'])
                )
                label = tbl['rows'][0]['cells'][0]['v'][:20] if tbl['rows'] else ''
                table_details.append(
                    f"    {tbl['rows_count']}r x {tbl['cols_count']}c "
                    f"merges={merges} [{label}]"
                )

        print(f"  Slide {sl['slide_no']}: {sl['shape_count']}shapes | {', '.join(parts)}",
              file=sys.stderr)
        for td in table_details:
            print(td, file=sys.stderr)


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='pptx에서 shapes_manifest.json 생성')
    parser.add_argument('pptx', help='레퍼런스 pptx 경로')
    parser.add_argument('output', nargs='?', default='shapes_manifest.json',
                        help='출력 JSON 경로 (기본: shapes_manifest.json)')
    parser.add_argument('--blank-layout-index', type=int, default=6,
                        help='빈 슬라이드 레이아웃 인덱스 (기본: 6)')
    parser.add_argument('--cell-margin-emu', type=int, default=36000,
                        help='기본 셀 마진 EMU (기본: 36000)')
    args = parser.parse_args()

    manifest = extract_all(args.pptx,
                           blank_layout_index=args.blank_layout_index,
                           cell_margin_emu=args.cell_margin_emu)
    result = json.dumps(manifest, ensure_ascii=False, indent=2)

    with open(args.output, 'w', encoding='utf-8') as f:
        f.write(result)
    print(f"→ {args.output} ({manifest['total_shapes']}shapes, {manifest['slide_count']}slides)")

    print_summary(manifest)
