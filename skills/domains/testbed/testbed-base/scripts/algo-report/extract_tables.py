"""레퍼런스 pptx에서 테이블만 일괄 추출하여 tables_detail.json 생성.

각 테이블의 셀 값과 원시 스타일 정보(볼드, 폰트크기, 배경색)를
있는 그대로 추출한다. H/V 판정은 하지 않는다 — 그 판단은 AI 의미 분석의 몫.

Usage:
    python extract_tables.py <input.pptx> [output.json]

출력 구조:
    [
      {
        "shape_id": "s3_shape2",
        "slide_no": 3,
        "rows_count": 7,
        "cols_count": 3,
        "col_widths_cm": [4.5, 6.2, 5.8],
        "rows": [
          {
            "row_idx": 0,
            "h_cm": 0.8,
            "cells": [
              {
                "v": "투자자 성향 구분",
                "bold": true,
                "font_size": 10.0,
                "fill_color": "#D9D9D9",
                "merge": {"cs": 2}
              }
            ]
          }
        ]
      }
    ]
"""
import sys, json
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.oxml.ns import qn


def _fill_color(tc):
    """셀의 배경 fill 색상을 hex 문자열로 반환. 없으면 None."""
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        return None
    fill = tcPr.find(qn('a:solidFill'))
    if fill is None:
        return None
    srgb = fill.find(qn('a:srgbClr'))
    if srgb is not None:
        return f"#{srgb.get('val', '000000')}"
    scheme = fill.find(qn('a:schemeClr'))
    if scheme is not None:
        val = scheme.get('val', '')
        lum = scheme.find(qn('a:lumMod'))
        lum_val = lum.get('val') if lum is not None else None
        return f"scheme:{val}" + (f":lum{lum_val}" if lum_val else "")
    return None


def _cell_bold(tc):
    """셀의 첫 run이 볼드인지 반환."""
    txBody = tc.find(qn('a:txBody'))
    if txBody is None:
        return None
    for p in txBody.findall(qn('a:p')):
        for r in p.findall(qn('a:r')):
            rPr = r.find(qn('a:rPr'))
            if rPr is not None:
                b = rPr.get('b')
                if b is not None:
                    return b == '1'
            return None
    return None


def _cell_font_size(tc):
    """셀의 첫 run 폰트 크기 (pt). 없으면 None."""
    txBody = tc.find(qn('a:txBody'))
    if txBody is None:
        return None
    for p in txBody.findall(qn('a:p')):
        for r in p.findall(qn('a:r')):
            rPr = r.find(qn('a:rPr'))
            if rPr is not None:
                sz = rPr.get('sz')
                if sz:
                    return round(int(sz) / 100, 1)
            return None
    return None


def _cell_merge(tc):
    """셀의 병합 정보."""
    merge = {}
    rs = tc.get('rowSpan')
    cs = tc.get('gridSpan')
    vm = tc.get('vMerge')
    hm = tc.get('hMerge')
    if rs:
        merge["rs"] = int(rs)
    if cs:
        merge["cs"] = int(cs)
    if vm:
        merge["vMerge"] = True
    if hm:
        merge["hMerge"] = True
    return merge or None


def extract_tables(pptx_path):
    """pptx에서 모든 테이블의 상세 정보를 추출."""
    prs = Presentation(pptx_path)
    tables = []

    for si, slide in enumerate(prs.slides):
        shape_counter = 0
        for shape in slide.shapes:
            shape_counter += 1
            if not shape.has_table:
                continue

            shape_id = f"s{si + 1}_shape{shape_counter}"
            tbl = shape.table
            tbl_el = tbl._tbl

            # 열 너비
            grid = tbl_el.find(qn('a:tblGrid'))
            col_widths_cm = []
            for gc in grid.findall(qn('a:gridCol')):
                col_widths_cm.append(round(int(gc.get('w')) / 360000, 2))

            rows = []
            for ri, tr in enumerate(tbl_el.findall(qn('a:tr'))):
                h_cm = round(int(tr.get('h')) / 360000, 2)
                cells = []
                for ci, tc in enumerate(tr.findall(qn('a:tc'))):
                    cell = {"v": tbl.cell(ri, ci).text.strip()}
                    bold = _cell_bold(tc)
                    if bold is not None:
                        cell["bold"] = bold
                    font_size = _cell_font_size(tc)
                    if font_size is not None:
                        cell["font_size"] = font_size
                    fill = _fill_color(tc)
                    if fill:
                        cell["fill_color"] = fill
                    merge = _cell_merge(tc)
                    if merge:
                        cell["merge"] = merge
                    cells.append(cell)
                rows.append({"row_idx": ri, "h_cm": h_cm, "cells": cells})

            tables.append({
                "shape_id": shape_id,
                "slide_no": si + 1,
                "rows_count": len(tbl.rows),
                "cols_count": len(tbl.columns),
                "col_widths_cm": col_widths_cm,
                "rows": rows,
            })

    return tables


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='pptx에서 테이블 상세 추출')
    parser.add_argument('pptx', help='레퍼런스 pptx 경로')
    parser.add_argument('output', nargs='?', default='tables_detail.json',
                        help='출력 JSON 경로 (기본: tables_detail.json)')
    args = parser.parse_args()

    tables = extract_tables(args.pptx)
    with open(args.output, 'w', encoding='utf-8') as f:
        json.dump(tables, f, ensure_ascii=False, indent=2)
    print(f'{len(tables)}개 테이블 → {args.output}')
    for t in tables:
        print(f'  {t["shape_id"]}: {t["rows_count"]}r x {t["cols_count"]}c (slide {t["slide_no"]})')
